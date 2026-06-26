"""Build the Strategy-D slide deck (PowerPoint, ~19 slides, editable diagrams).

Diagrams are native PowerPoint shapes (rounded rectangles + connectors with
arrowheads) so the supervisor can edit them; the three analytical charts are the
shared matplotlib PNGs. Run:

    py -3.11 presentation/strategy_d/make_pptx.py

Output: presentation/strategy_d/strategy_d_slides.pptx
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import charts  # noqa: E402
from flickr_results_data import (  # noqa: E402
    CHEN_PAIRS_F1, CITIES, DATASET_STATS, OURS_PAIRS_F1, PAIRS_EXAMPLE,
    POINTER_PAIRS_F1, PUB_CITIES, PUBLISHED_PAIRS_F1,
)

ASSETS = charts.make_charts(os.path.join(HERE, "assets"))
OUTPUT = os.path.join(HERE, "strategy_d_slides.pptx")

# ---- palette ----
NAVY = RGBColor(0x1B, 0x2A, 0x4E)
INDIGO = RGBColor(0x3F, 0x51, 0xB5)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
ORANGE = RGBColor(0xE0, 0x7A, 0x1F)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
GREY_MED = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_G = RGBColor(0xE3, 0xF2, 0xFD)
BG_C = RGBColor(0xFF, 0xF3, 0xE0)
BG_U = RGBColor(0xF3, 0xE5, 0xF5)
BG_S = RGBColor(0xE8, 0xF5, 0xE9)
BG_H = RGBColor(0xFF, 0xEB, 0xEE)
BG_LIGHT = RGBColor(0xEC, 0xEE, 0xF3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


# ===========================================================================
# Primitives
# ===========================================================================
def slide():
    return prs.slides.add_slide(BLANK)


def _no_shadow(shp):
    shp.shadow.inherit = False


def rect(s, x, y, w, h, fill, line=None, line_w=1.5, rounded=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    _no_shadow(shp)
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """runs = list of paragraphs; each paragraph = list of (txt, size, color, bold, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for pi, para in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold, italic) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = "Calibri"
    return tb


def box(s, x, y, w, h, lines, fill, line, anchor=MSO_ANCHOR.MIDDLE):
    """A diagram box: lines = list of (txt, size, color, bold)."""
    shp = rect(s, x, y, w, h, fill, line, line_w=1.75)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    for i, (txt, size, col, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = col
        r.font.name = "Calibri"
    return shp


def arrow(s, x1, y1, x2, y2, color=NAVY, width=2.0):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    _no_shadow(conn)
    return conn


# ---- slide chrome ----
def header(s, title, kicker=None):
    rect(s, 0, 0, SW, 1.15, NAVY, line=None, rounded=False)
    rect(s, 0, 1.15, SW, 0.06, TEAL, line=None, rounded=False)
    text(s, 0.55, 0.18, SW - 1.1, 0.85,
         [[(title, 28, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(s, 0.55, 0.05, SW - 1.1, 0.3, [[(kicker, 11, RGBColor(0xBF, 0xC7, 0xD9), False, True)]])


def footer(s, n):
    text(s, SW - 1.4, SH - 0.45, 1.1, 0.3, [[(f"{n}", 11, GREY_MED, False, False)]], align=PP_ALIGN.RIGHT)
    text(s, 0.55, SH - 0.45, 8, 0.3,
         [[("Strategy D — Itinerary recommendation on Flickr", 9, GREY_MED, False, True)]])


def bullets(s, x, y, w, h, items, size=17, gap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10 if gap else 4)
        p.space_before = Pt(2)
        # item = list of (txt, color, bold) runs, or a string
        if isinstance(item, str):
            item = [(item, GREY, False)]
        bull = p.add_run()
        bull.text = "▸  "
        bull.font.size = Pt(size)
        bull.font.color.rgb = TEAL
        bull.font.bold = True
        for (txt, color, bold) in item:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = "Calibri"
    return tb


def image(s, path, x, y, w):
    return s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


def table(s, rows, x, y, w, h, header_fill=NAVY, first_col_left=True, highlight_row=None,
          font=12):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for ci in range(nc):
        for ri in range(nr):
            cell = gt.cell(ri, ci)
            cell.margin_top = Pt(3)
            cell.margin_bottom = Pt(3)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (ci == 0 and first_col_left) else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(rows[ri][ci])
            r.font.size = Pt(font)
            r.font.name = "Calibri"
            if ri == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                r.font.color.rgb = WHITE
                r.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else BG_LIGHT
                r.font.color.rgb = GREY
                if highlight_row is not None and ri == highlight_row:
                    cell.fill.fore_color.rgb = TEAL
                    r.font.color.rgb = WHITE
                    r.font.bold = True
    return gt


# ===========================================================================
# Slides
# ===========================================================================
def s_title():
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY, line=None, rounded=False)
    rect(s, 0, 3.05, SW, 0.05, TEAL, line=None, rounded=False)
    text(s, 1.0, 1.4, SW - 2, 1.5,
         [[("Literature-Comparable", 40, WHITE, True, False)],
          [("Itinerary Recommendation on Flickr", 40, WHITE, True, False)]],
         align=PP_ALIGN.CENTER)
    text(s, 1.0, 3.25, SW - 2, 0.7,
         [[("Strategy D — putting our pairs-F1 on the same scale as the literature", 20,
            RGBColor(0x9F, 0xC0, 0xC0), False, True)]], align=PP_ALIGN.CENTER)
    text(s, 1.0, 4.5, SW - 2, 1.6,
         [[("From ", 17, WHITE, False, False), ("non-comparable ", 17, ORANGE, True, False),
           ("Foursquare-NYC scores  →  pairs-F1 beside Chen 2016, DeepTrip,", 17, WHITE, False, False)],
          [("SelfTrip and AR-Trip on the Flickr photo-trajectory benchmark", 17, WHITE, False, False)]],
         align=PP_ALIGN.CENTER)
    text(s, 1.0, 6.4, SW - 2, 0.5,
         [[("Master's thesis — Phase 2 · github.com/6ym6n/PFE_IMPLEMTATION", 13,
            RGBColor(0xBF, 0xC7, 0xD9), False, False)]], align=PP_ALIGN.CENTER)


def s_where(n):
    s = slide()
    header(s, "Where Strategy D fits", "context")
    text(s, 0.55, 1.45, SW - 1.1, 0.9,
         [[("The thesis has ", 17, GREY, False, False), ("two formally distinct tasks", 17, NAVY, True, False),
           (", each on its own field-standard benchmark. Strategy D is the itinerary task, on Flickr.",
            17, GREY, False, False)]])
    y = 3.2
    bw, bh = 3.95, 2.4
    data = [
        (0.55, BG_G, INDIGO, "Phase 1", "Next-POI prediction",
         ["Foursquare NYC", "GCN + GRU + user", "HR@k / NDCG / MRR"]),
        (4.7, BG_H, ORANGE, "Phase 2A / 2B", "Itinerary on NYC",
         ["frozen rollout / pointer", "pairs-F1 0.26–0.29", "NOT comparable"]),
        (8.85, BG_S, TEAL, "Phase 2D — THIS WORK", "Itinerary on Flickr",
         ["Chen-2016 protocol", "pairs-F1 0.23–0.59", "literature-comparable"]),
    ]
    for x, fill, stroke, t1, t2, subs in data:
        shp = box(s, x, y, bw, bh,
                  [(t1, 15, stroke, True), (t2, 15, stroke, True), ("", 5, stroke, False)]
                  + [(ln, 12.5, GREY, False) for ln in subs], fill, stroke)
    for i in range(2):
        x1 = 0.55 + bw + i * 4.15
        arrow(s, x1, y + bh / 2, x1 + 0.6, y + bh / 2)
    text(s, 0.55, 6.2, SW - 1.1, 0.7,
         [[("Using a different dataset per task is ", 14, GREY, False, False),
           ("correct and standard", 14, NAVY, True, False),
           (" — the only error would be to cross-compare the two.", 14, GREY, False, False)]])
    footer(s, n)


def s_problem(n):
    s = slide()
    header(s, "The problem", "motivation")
    text(s, 0.55, 1.4, SW - 1.1, 1.0,
         [[("Our Phase-2 itinerary model scored ", 18, GREY, False, False),
           ("pairs-F1 ≈ 0.26–0.29 on Foursquare NYC", 18, RED, True, False),
           (".  The literature reports ", 18, GREY, False, False),
           ("0.6–0.8", 18, TEAL, True, False),
           (".  Are we bad?", 18, GREY, False, False)]])
    image(s, ASSETS["scale"], 1.4, 2.5, 10.5)
    text(s, 0.55, 6.55, SW - 1.1, 0.7,
         [[("No — not a quality gap, a ", 18, GREY, False, False),
           ("comparability gap", 18, NAVY, True, False),
           (".  Different benchmark → different scale.", 18, GREY, False, False)]],
         align=PP_ALIGN.CENTER)
    footer(s, n)


def s_insight(n):
    s = slide()
    header(s, "What makes two numbers comparable", "the insight")
    text(s, 0.55, 1.45, SW - 1.1, 0.6,
         [[("Two pairs-F1 values compare ", 18, GREY, False, False),
           ("only if all three match:", 18, NAVY, True, False)]])
    cards = [
        (0.7, BG_G, INDIGO, "1.  DATASET", ["same trajectories,", "same POI-vocabulary size",
                                            "(30 POIs vs 5,000 →", "different chance levels)"]),
        (4.85, BG_C, ORANGE, "2.  PROTOCOL", ["same split,", "same query,", "same length filter"]),
        (9.0, BG_S, TEAL, "3.  METRIC", ["the exact same", "pairs-F1 definition", "(pinned by a unit test)"]),
    ]
    for x, fill, stroke, title, subs in cards:
        box(s, x, 2.4, 3.6, 2.7,
            [(title, 18, stroke, True), ("", 8, stroke, False)] + [(ln, 13.5, GREY, False) for ln in subs],
            fill, stroke)
    text(s, 0.55, 5.5, SW - 1.1, 1.2,
         [[("Foursquare NYC differs on ", 17, GREY, False, False), ("all three", 17, RED, True, False),
           (", so its 0.26–0.29 is incomparable to the literature's 0.6–0.8.", 17, GREY, False, False)],
          [("The gap is a data/protocol artefact, not model quality (Krichene & Rendle, KDD 2020).",
            14, GREY_MED, False, True)]])
    footer(s, n)


def s_datasets(n):
    s = slide()
    header(s, "The Flickr datasets", "data")
    text(s, 0.55, 1.4, SW - 1.1, 1.2,
         [[("Five cities mined from geo-tagged Flickr photos. We use the ", 15, GREY, False, False),
           ("canonical traj-/poi- CSVs = Chen 2016's own files", 15, NAVY, True, False),
           (", so our trajectories are identical to the literature's (zero preprocessing ambiguity).",
            15, GREY, False, False)]])
    # trajectory mini-diagram
    y = 2.7
    steps = [("Flickr photos", "(lat, lon, time)", BG_G, INDIGO),
             ("POI visits", "merge same-POI run", BG_C, ORANGE),
             ("order by time", "one seqID = trip", BG_U, INDIGO),
             ("trajectory", "p1→p2→p3 (loop-free)", BG_S, TEAL)]
    bw = 2.7
    for i, (t1, t2, fill, stroke) in enumerate(steps):
        x = 0.55 + i * 3.05
        box(s, x, y, bw, 1.2, [(t1, 13, stroke, True), (t2, 10.5, GREY, False)], fill, stroke)
        if i < 3:
            arrow(s, x + bw, y + 0.6, x + bw + 0.35, y + 0.6)
    # stats table
    rows = [["City", "#POIs", "#users", "#traj", "#traj≥3 (eval)"]]
    for c in CITIES:
        p, u, t, e = DATASET_STATS[c]
        rows.append([c, p, f"{u:,}", f"{t:,}", e])
    table(s, rows, 2.4, 4.35, 8.5, 2.4, font=13)
    text(s, 0.55, 6.95, SW - 1.1, 0.4,
         [[("The length≥3 column is the evaluated set; Chen 2016's totals match #traj exactly.",
            12, GREY_MED, False, True)]], align=PP_ALIGN.CENTER)
    footer(s, n)


def s_task(n):
    s = slide()
    header(s, "The task, precisely", "problem formulation")
    text(s, 0.55, 1.45, SW - 1.1, 0.9,
         [[("A ", 18, GREY, False, False), ("query", 18, NAVY, True, False),
           (" gives the start POI, the end POI and the length K. The model recovers the ", 18, GREY, False, False),
           ("ordered middle", 18, ORANGE, True, False), (".", 18, GREY, False, False)]])
    y = 3.1
    box(s, 0.55, y, 3.6, 1.6, [("QUERY", 15, INDIGO, True), ("start s · end e · length K", 13, GREY, False),
                               ("e.g. (Castle, Museum, 5)", 11, GREY_MED, False)], BG_U, INDIGO)
    arrow(s, 4.15, y + 0.8, 4.8, y + 0.8)
    box(s, 4.8, y + 0.15, 2.6, 1.3, [("recommender", 15, TEAL, True)], BG_S, TEAL)
    arrow(s, 7.4, y + 0.8, 8.05, y + 0.8)
    labels = [("s", BG_U), ("?", WHITE), ("?", WHITE), ("?", WHITE), ("e", BG_U)]
    for i, (lab, fl) in enumerate(labels):
        x = 8.1 + i * 1.0
        box(s, x, y + 0.25, 0.85, 1.1, [(lab, 18, ORANGE, True)], fl, ORANGE)
    text(s, 8.1, y + 1.5, 5, 0.5, [[("ordered route of K POIs, loop-free, ends at e", 12, GREY_MED, False, True)]])
    text(s, 0.55, 5.6, SW - 1.1, 1.0,
         [[("Standard PersTour / Chen-2016 convention: ", 16, GREY, False, False),
           ("recover the visited sequence given its endpoints and length.", 16, NAVY, True, False)],
          [("Start and end are given, so only the middle ordering is at stake.", 14, GREY_MED, False, True)]])
    footer(s, n)


def s_metric(n):
    s = slide()
    header(s, "The metric: pairs-F1 (Chen 2016)", "the key idea")
    text(s, 0.55, 1.35, SW - 1.1, 0.7,
         [[("pairs-F1 = F1 over ", 17, GREY, False, False), ("ordered pairs", 17, NAVY, True, False),
           (". A pair (a,b) counts only if a is before b in ", 17, GREY, False, False),
           ("both", 17, ORANGE, True, False), (" routes.", 17, GREY, False, False)]])
    ex = PAIRS_EXAMPLE
    # two sequences
    def seqrow(label, seq, y, color, fill):
        text(s, 0.55, y, 2.6, 0.6, [[(label, 14, color, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
        for i, p in enumerate(seq):
            box(s, 3.2 + i * 1.0, y, 0.85, 0.6, [(p, 16, color, True)], fill, color)
    seqrow("ground truth Y:", ex["truth"], 2.35, TEAL, BG_S)
    seqrow("prediction Ŷ:", ex["pred"], 3.15, ORANGE, BG_H)
    text(s, 0.55, 4.05, 12, 1.4,
         [[("Y pairs:  ", 13, NAVY, True, False), ("  ".join(ex["truth_pairs"]), 13, TEAL, False, False)],
          [("Ŷ pairs:  ", 13, NAVY, True, False), ("  ".join(ex["pred_pairs"]), 13, ORANGE, False, False)],
          [("shared, same order:  ", 13, NAVY, True, False),
           ("  ".join(ex["shared"]) + "   →  5 of 6", 13, NAVY, True, False)],
          [("(B,C) is in Y but Ŷ has (C,B) — wrong order, doesn't count.", 11.5, GREY_MED, False, True)]])
    box(s, 0.55, 5.9, 12.2, 0.9,
        [("precision = 5/6,  recall = 5/6   →   pairs-F1 = 2·P·R/(P+R) = 5/6 ≈ 0.833", 18, NAVY, True)],
        BG_LIGHT, NAVY)
    footer(s, n)


def s_protocol(n):
    s = slide()
    header(s, "The protocol: leave-one-trajectory-out CV", "evaluation")
    text(s, 0.55, 1.4, SW - 1.1, 0.9,
         [[("Tiny datasets (47–634 trajectories), so the field uses ", 16, GREY, False, False),
           ("leave-one-out cross-validation", 16, NAVY, True, False),
           (".", 16, GREY, False, False)]])
    # column of trajectories
    names = ["t₁", "t₂", "t₃  ← held out (TEST)", "t₄", "t₅", "t₆"]
    for i, nm in enumerate(names):
        is_test = i == 2
        box(s, 0.55, 2.5 + i * 0.62, 3.4, 0.5, [(nm, 12, ORANGE if is_test else TEAL, is_test)],
            BG_H if is_test else BG_S, ORANGE if is_test else TEAL)
    arrow(s, 4.05, 3.7, 5.0, 3.7)
    box(s, 5.0, 3.1, 3.5, 1.2, [("train on the other N−1", 14, TEAL, True),
                                ("(length≥3 trajectories)", 11, GREY, False)], BG_S, TEAL)
    arrow(s, 8.5, 3.7, 9.4, 3.7)
    box(s, 9.4, 3.1, 3.4, 1.2, [("predict t₃ from (s,e,K)", 13, ORANGE, True),
                                ("score F1 & pairs-F1", 13, ORANGE, True)], BG_H, ORANGE)
    text(s, 0.55, 6.4, SW - 1.1, 0.7,
         [[("Repeat holding out every trajectory once → report the mean per city. ", 16, NAVY, True, False),
           ("No leakage: the held-out trajectory never touches training.", 14, GREY_MED, False, True)]])
    footer(s, n)


def s_methods(n):
    s = slide()
    header(s, "What we built", "methods")
    # classical band
    text(s, 0.55, 1.35, 6, 0.4, [[("Classical baselines (counting, CPU)", 16, NAVY, True, False)]])
    cb = [("training\ntrajectories", BG_S, TEAL), ("popularity /\nMarkov P(j|i)", BG_C, ORANGE),
          ("greedy / beam\ndecode (loop-free)", BG_H, ORANGE)]
    for i, (t, fill, stroke) in enumerate(cb):
        x = 0.55 + i * 4.2
        lines = [(ln, 13, stroke if k == 0 else GREY, k == 0) for k, ln in enumerate(t.split("\n"))]
        box(s, x, 1.85, 3.6, 1.15, lines, fill, stroke)
        if i < 2:
            arrow(s, x + 3.6, 2.42, x + 3.6 + 0.6, 2.42)
    # mapping to Chen
    text(s, 0.55, 3.1, 12, 0.4,
         [[("= Chen 2016's Random · PoiPopularity · Markov · MarkovPath", 12.5, GREY_MED, False, True)]])
    # neural band
    text(s, 0.55, 3.6, 9, 0.4, [[("Learned model: GCN encoder + GRU pointer (GPU, no torch_geometric)",
                                  16, NAVY, True, False)]])
    arch = [("POI graph\nkNN ∪ co-visit", BG_G, INDIGO), ("2-layer GCN", BG_G, INDIGO),
            ("POI features H", BG_G, INDIGO), ("GRU pointer\n⟨state, H_v⟩", BG_S, TEAL),
            ("next POI\nmask+reserve e", BG_H, ORANGE)]
    bw = 2.3
    for i, (t, fill, stroke) in enumerate(arch):
        x = 0.55 + i * 2.55
        lines = [(ln, 12, stroke if k == 0 else GREY, k == 0) for k, ln in enumerate(t.split("\n"))]
        box(s, x, 4.15, bw, 1.15, lines, fill, stroke)
        if i < 4:
            arrow(s, x + bw, 4.72, x + bw + 0.25, 4.72)
    bullets(s, 0.55, 5.7, 12.2, 1.4, [
        [("A fresh model is trained on every leave-one-out fold; the query (s,e) seeds the decoder.",
          GREY, False)],
        [("Levers to close the gap: ", NAVY, True), ("Markov transition prior at decode, user embedding, early stopping.",
          GREY, False)],
    ], size=15)
    footer(s, n)


def s_validation(n):
    s = slide()
    header(s, "Did we do it right? Reproducing Chen 2016", "validation")
    text(s, 0.55, 1.35, 6.6, 1.8,
         [[("A ", 15, GREY, False, False), ("Random", 15, NAVY, True, False),
           (" baseline has no modelling choices — it can only match Chen's Random if our data, "
            "protocol and metric are identical.", 15, GREY, False, False)],
          [("", 6, GREY, False, False)],
          [("They match within noise on every city; PoiPopularity matches closely too.", 15, NAVY, True, False)]])
    rows = [["pairs-F1", "Tor", "Osa", "Gla", "Edi", "Mel"]]
    for m in ("Random", "PoiPopularity"):
        rows.append([f"{m} — ours"] + [f"{OURS_PAIRS_F1[m][c]:.3f}" for c in CITIES])
        rows.append([f"{m} — Chen"] + [f"{CHEN_PAIRS_F1[m][c]:.3f}" for c in CITIES])
    table(s, rows, 0.55, 4.0, 6.6, 2.4, font=12)
    image(s, ASSETS["validation"], 7.7, 1.5, 5.2)
    footer(s, n)


def s_results(n):
    s = slide()
    header(s, "Results — on the literature's scale", "results")
    image(s, ASSETS["results"], 0.7, 1.45, 9.0)
    box(s, 9.95, 1.7, 3.0, 4.2,
        [("Headline", 17, NAVY, True), ("", 6, NAVY, False),
         ("our classical", 14, GREY, True), ("0.23 – 0.59", 22, TEAL, True), ("", 5, NAVY, False),
         ("our pointer", 14, GREY, True), ("0.31 – 0.49", 22, ORANGE, True), ("", 5, NAVY, False),
         ("vs NYC", 14, GREY, True), ("0.26 – 0.29", 18, RED, True), ("", 5, NAVY, False),
         ("published", 14, GREY, True), ("0.26 – 0.85", 16, NAVY, True)],
        BG_LIGHT, NAVY)
    text(s, 0.7, 6.5, 9, 0.6,
         [[("Every one of our methods is on the published 0–1 scale — the thesis goal is met by the "
            "classical baselines alone.", 15, NAVY, True, False)]])
    footer(s, n)


def s_published(n):
    s = slide()
    header(s, "The published comparison", "literature")
    text(s, 0.55, 1.35, SW - 1.1, 0.5,
         [[("Directly-comparable published pairs-F1 (same protocol, 0–1 scale):", 16, GREY, False, False)]])
    rows = [["method (published)"] + PUB_CITIES]
    for m, (vals, fam, yr) in PUBLISHED_PAIRS_F1.items():
        rows.append([m] + [f"{v:.3f}" for v in vals])
    table(s, rows, 1.6, 2.1, 10.1, 3.4, font=14)
    text(s, 0.55, 5.9, SW - 1.1, 1.0,
         [[("No single SOTA across all cities: ", 15, GREY, False, False),
           ("AR-Trip leads on Toronto/Glasgow/Edinburgh, SelfTrip on Osaka (0.851).", 15, NAVY, True, False)],
          [("Classical methods (Chen 2016) ≈ 0.43–0.55; neural methods (2019–2024) climb to ≈ 0.66–0.85.",
            13, GREY_MED, False, True)]])
    footer(s, n)


def s_learned(n):
    s = slide()
    header(s, "The learned model — measured & honest", "the finding")
    rows = [["", "Tor", "Osa", "Gla", "Edi", "Mel"]]
    rows.append(["Pointer pairs-F1"] + [f"{POINTER_PAIRS_F1['Pointer (beam)'][c]:.3f}" for c in CITIES])
    rows.append(["Markov (ours)"] + [f"{OURS_PAIRS_F1['Markov'][c]:.3f}" for c in CITIES])
    rows.append(["Random (ours)"] + [f"{OURS_PAIRS_F1['Random'][c]:.3f}" for c in CITIES])
    table(s, rows, 0.55, 1.45, 7.0, 1.9, font=13, highlight_row=1)
    bullets(s, 0.55, 3.7, 12.3, 3.2, [
        [("On the scale, above Random: ", NAVY, True),
         ("pointer reaches 0.31–0.49, with point-F1 ~0.70 (it picks the right POIs).", GREY, False)],
        [("But its ordering trails Markov on every city", RED, True),
         (" and is far below the neural SOTA (0.66–0.85).", GREY, False)],
        [("The lesson (same as Strategy B on NYC): ", NAVY, True),
         ("naively training a pointer on tiny per-fold data doesn't beat a strong simple model.", GREY, False)],
        [("SOTA methods use machinery we omit: ", GREY, False),
         ("self-supervised pre-training, trajectory augmentation, adversarial training.", GREY, False)],
    ], size=16)
    footer(s, n)


def s_levers(n):
    s = slide()
    header(s, "Closing the gap — implemented levers", "next iteration")
    cards = [
        (0.7, BG_C, ORANGE, "Markov prior", ["blend the fold's log P(j|i)", "into the pointer logits",
                                             "at decode (no retraining)", "→ targets the weak ordering"]),
        (4.85, BG_U, INDIGO, "User embedding", ["per-user vector added to", "the query + decoder state"]),
        (9.0, BG_S, TEAL, "Early stopping", ["hold out a val slice of", "the train fold; stop on", "val loss"]),
    ]
    for x, fill, stroke, title, subs in cards:
        box(s, x, 1.7, 3.6, 2.5,
            [(title, 18, stroke, True), ("", 8, stroke, False)] + [(ln, 13, GREY, False) for ln in subs],
            fill, stroke)
    box(s, 0.7, 4.6, 11.9, 1.0,
        [("Verified: Osaka pairs-F1  0.42 (pure) → 0.44 (+Markov+user) — edges past the Markov baseline.",
          17, NAVY, True)], BG_LIGHT, TEAL)
    text(s, 0.7, 5.9, 11.9, 0.9,
         [[("All opt-in in ", 14, GREY, False, False), ("PointerConfig", 14, NAVY, True, False),
           ("; the Colab notebook runs the pure pointer and the enhanced variant side by side. "
            "Sweep the prior weight / epochs for a higher number.", 14, GREY, False, False)]])
    footer(s, n)


def s_value(n):
    s = slide()
    header(s, "What this gives the thesis", "contribution")
    bullets(s, 0.7, 1.6, 12.0, 5.0, [
        [("Literature-comparable numbers", NAVY, True),
         (" — pairs-F1 on the same 0.3–0.85 scale as Chen 2016 / DeepTrip / SelfTrip / AR-Trip.", GREY, False)],
        [("A validated, faithful harness", NAVY, True),
         (" — Random reproduces Chen 2016 within noise (the protocol-faithfulness proof).", GREY, False)],
        [("An honest baseline", NAVY, True),
         (" — classical methods are strong; the light learned model reaches the scale but not SOTA, "
          "and we say so (with the levers to improve it).", GREY, False)],
        [("Fully reproducible", NAVY, True),
         (" — self-contained loader + harness + Colab notebook; 79 passing tests; no heavy deps.", GREY, False)],
        [("Resolves the NYC puzzle", NAVY, True),
         (" — the low Foursquare-NYC pairs-F1 was a data/protocol artefact, not a bug.", GREY, False)],
    ], size=18)
    footer(s, n)


def s_summary(n):
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY, line=None, rounded=False)
    rect(s, 0, 1.7, SW, 0.05, TEAL, line=None, rounded=False)
    text(s, 0.8, 0.6, SW - 1.6, 1.0, [[("Summary", 34, WHITE, True, False)]])
    bullets_white = [
        "Same itinerary task, the field's benchmark (Flickr) + protocol (Chen 2016, leave-one-out, endpoints given, pairs-F1).",
        "Random reproduces Chen → faithful.  Classical 0.23–0.59, pointer 0.31–0.49 — on the published 0.26–0.85 scale.",
        "Honest finding: the light learned pointer trails the Markov baseline and SOTA; levers implemented to close the gap.",
    ]
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(SW - 1.6), Inches(2.6))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets_white):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        b = p.add_run(); b.text = "▸  "; b.font.size = Pt(18); b.font.color.rgb = TEAL; b.font.bold = True
        r = p.add_run(); r.text = item; r.font.size = Pt(18); r.font.color.rgb = WHITE; r.font.name = "Calibri"
    text(s, 0.8, 5.1, SW - 1.6, 0.5, [[("Next steps", 20, ORANGE, True, False)]])
    tb2 = s.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(SW - 1.6), Inches(1.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, item in enumerate([
        "Run the enhanced pointer on Colab (sweep Markov-prior weight / epochs) for a stronger learned number.",
        "Fold these results into the full thesis defence deck (next session).",
    ]):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.space_after = Pt(8)
        b = p.add_run(); b.text = "•  "; b.font.size = Pt(15); b.font.color.rgb = ORANGE
        r = p.add_run(); r.text = item; r.font.size = Pt(15); r.font.color.rgb = RGBColor(0xD8, 0xDE, 0xEA); r.font.name = "Calibri"


def s_appendix(n):
    s = slide()
    header(s, "Appendix — reproducibility & references", "backup")
    text(s, 0.55, 1.35, 6.2, 0.4, [[("Code (src/flickr/)", 16, NAVY, True, False)]])
    bullets(s, 0.55, 1.8, 6.4, 3.5, [
        "data.py — loader (both formats) → FlickrCity",
        "evaluate.py — leave-one-out CV, F1 + pairs-F1",
        "baselines.py — Random/Popularity/Markov(Path)",
        "pointer.py — GCN+pointer + levers (Colab)",
        "published.py — curated literature numbers",
        "run_flickr.py — orchestration + tables + CLI",
        "tests/test_flickr.py — 79 cases incl. pairs-F1 vs Chen",
    ], size=13.5, gap=False)
    text(s, 7.1, 1.35, 5.7, 0.4, [[("Key references", 16, NAVY, True, False)]])
    bullets(s, 7.1, 1.8, 5.8, 4.0, [
        "Chen, Ong, Xie. Learning Points & Routes. CIKM 2016 (pairs-F1 + our data).",
        "Gao et al. DeepTrip. SIGSPATIAL 2019.",
        "Gao et al. SelfTrip. KBS 2022.",
        "Rashid et al. DeepAltTrip. TKDE 2021.",
        "Shu et al. AR-Trip. SIGIR 2024.",
        "Lim et al. Tour Rec & Trip Planning survey. KAIS 2019.",
        "Krichene & Rendle. Sampled Metrics. KDD 2020.",
    ], size=13.5, gap=False)
    footer(s, n)


def build():
    s_title()
    s_where(2)
    s_problem(3)
    s_insight(4)
    s_datasets(5)
    s_task(6)
    s_metric(7)
    s_protocol(8)
    s_methods(9)
    s_validation(10)
    s_results(11)
    s_published(12)
    s_learned(13)
    s_levers(14)
    s_value(15)
    s_summary(16)
    s_appendix(17)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
