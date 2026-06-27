"""Full-thesis defence deck (~21 slides).

Tells the whole arc: personalized context-aware itinerary recommendation, with the
next-POI model as the engine, a decoupled-vs-integrated itinerary experiment
(Strategy A vs B), and a literature-comparable Flickr benchmark (Strategy D) — and
positions it against Halder's survey / DLIR. Self-contained native-shape diagrams +
the shared charts. Run:

    py -3.11 presentation/strategy_d/make_thesis_pptx.py

Output: presentation/strategy_d/thesis_defense_slides.pptx
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import charts  # noqa: E402
from flickr_results_data import (  # noqa: E402
    AVSB_NYC, CHEN_PAIRS_F1, CITIES, DATASET_STATS, OURS_PAIRS_F1, PHASE1_METRICS,
    PHASE1_TIER, POINTER_PAIRS_F1, PUBLISHED_PAIRS_F1,
)

ASSETS = charts.make_charts(os.path.join(HERE, "assets"))
OUTPUT = os.path.join(HERE, "thesis_defense_slides.pptx")

NAVY = RGBColor(0x1B, 0x2A, 0x4E)
INDIGO = RGBColor(0x3F, 0x51, 0xB5)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
ORANGE = RGBColor(0xE0, 0x7A, 0x1F)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
GREY_MED = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_G = RGBColor(0xE3, 0xF2, 0xFD)
BG_C = RGBColor(0xFF, 0xF3, 0xE0)
BG_U = RGBColor(0xF3, 0xE5, 0xF5)
BG_S = RGBColor(0xE8, 0xF5, 0xE9)
BG_H = RGBColor(0xFF, 0xEB, 0xEE)
BG_LIGHT = RGBColor(0xEC, 0xEE, 0xF3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


# ---- primitives (shared with the Strategy-D deck) ----
def slide():
    return prs.slides.add_slide(BLANK)


def _ns(shp):
    shp.shadow.inherit = False


def rect(s, x, y, w, h, fill, line=None, line_w=1.5, rounded=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    _ns(shp)
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for pi, para in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic; r.font.name = "Calibri"
    return tb


def box(s, x, y, w, h, lines, fill, line, anchor=MSO_ANCHOR.MIDDLE):
    shp = rect(s, x, y, w, h, fill, line, line_w=1.75)
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, (txt, size, col, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
    return shp


def arrow(s, x1, y1, x2, y2, color=NAVY, width=2.0):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    _ns(conn)
    return conn


def header(s, title, kicker=None):
    rect(s, 0, 0, SW, 1.15, NAVY, line=None, rounded=False)
    rect(s, 0, 1.15, SW, 0.06, TEAL, line=None, rounded=False)
    text(s, 0.55, 0.18, SW - 1.1, 0.85, [[(title, 27, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(s, 0.55, 0.04, SW - 1.1, 0.3, [[(kicker, 11, RGBColor(0xBF, 0xC7, 0xD9), False, True)]])


def footer(s, n):
    text(s, SW - 1.4, SH - 0.45, 1.1, 0.3, [[(f"{n}", 11, GREY_MED, False, False)]], align=PP_ALIGN.RIGHT)
    text(s, 0.55, SH - 0.45, 10, 0.3,
         [[("Smart Visit Module — Personalized Tourist Itinerary Recommendation", 9, GREY_MED, False, True)]])


def bullets(s, x, y, w, h, items, size=17, gap=10):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(2)
        if isinstance(item, str):
            item = [(item, GREY, False)]
        b = p.add_run(); b.text = "▸  "; b.font.size = Pt(size); b.font.color.rgb = TEAL; b.font.bold = True
        for (txt, color, bold) in item:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = "Calibri"
    return tb


def image(s, path, x, y, w):
    return s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


def table(s, rows, x, y, w, h, header_fill=NAVY, highlight_row=None, font=12):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for ci in range(nc):
        for ri in range(nr):
            cell = gt.cell(ri, ci)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(rows[ri][ci]); r.font.size = Pt(font); r.font.name = "Calibri"
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                r.font.color.rgb = WHITE; r.font.bold = True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else BG_LIGHT
                r.font.color.rgb = GREY
                if highlight_row is not None and ri == highlight_row:
                    cell.fill.fore_color.rgb = TEAL; r.font.color.rgb = WHITE; r.font.bold = True
    return gt


# ===========================================================================
# Slides
# ===========================================================================
def s_title():
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY, line=None, rounded=False)
    rect(s, 0, 3.35, SW, 0.05, TEAL, line=None, rounded=False)
    text(s, 1.0, 1.5, SW - 2, 1.8,
         [[("Smart Visit Module for Personalized", 33, WHITE, True, False)],
          [("Tourist Itinerary Recommendation", 33, WHITE, True, False)]], align=PP_ALIGN.CENTER)
    text(s, 1.0, 3.55, SW - 2, 0.7,
         [[("Based on User Preferences and Contextual Data", 19, RGBColor(0x9F, 0xC0, 0xC0), False, True)]],
         align=PP_ALIGN.CENTER)
    text(s, 1.0, 5.0, SW - 2, 1.2,
         [[("A personalized, context-aware next-POI engine, decoded and trained into full", 15, WHITE, False, False)],
          [("itineraries, and benchmarked against the trip-recommendation literature", 15, WHITE, False, False)]],
         align=PP_ALIGN.CENTER)
    text(s, 1.0, 6.5, SW - 2, 0.5,
         [[("Master's thesis · [Your name] · Supervisor: [name] · [date]", 13, RGBColor(0xBF, 0xC7, 0xD9), False, False)]],
         align=PP_ALIGN.CENTER)


def s_motivation(n):
    s = slide()
    header(s, "Motivation", "the problem")
    bullets(s, 0.7, 1.6, 12.0, 3.0, [
        [("A tourist with limited time faces hundreds of POIs and wants a ", GREY, False),
         ("personalized day plan", NAVY, True), (" — an ", GREY, False),
         ("ordered route", NAVY, True), (" of places to visit, not just a single suggestion.", GREY, False)],
        [("A good plan depends on ", GREY, False), ("who they are", NAVY, True),
         (" (preferences) and ", GREY, False), ("context", NAVY, True),
         (" (where they are, how far apart places are, time of day).", GREY, False)],
        [("This is ", GREY, False), ("itinerary recommendation", NAVY, True),
         (" — a different problem from predicting the single next place (next-POI).", GREY, False)],
    ], size=18)
    box(s, 1.6, 4.9, 10.1, 1.4,
        [("Goal: a “Smart Visit Module” that recommends a personalized, context-aware itinerary —",
          18, NAVY, True),
         ("the right places, in the right order, for this user.", 18, NAVY, True)], BG_S, TEAL)
    footer(s, n)


def s_gap(n):
    s = slide()
    header(s, "Two sub-problems — and the gap", "background")
    box(s, 0.7, 1.7, 5.6, 1.9,
        [("Next-POI prediction", 17, INDIGO, True), ("predict the single NEXT place", 14, GREY, False),
         ("(a building block)", 12, GREY_MED, False)], BG_G, INDIGO)
    box(s, 7.0, 1.7, 5.6, 1.9,
        [("Itinerary recommendation", 17, TEAL, True), ("produce the WHOLE ordered route", 14, GREY, False),
         ("(the actual goal)", 12, GREY_MED, False)], BG_S, TEAL)
    arrow(s, 6.3, 2.65, 7.0, 2.65)
    box(s, 1.6, 4.4, 10.1, 1.9,
        [("The field's question (Halder 2024 survey; DLIR 2025):", 16, NAVY, True),
         ("should POI scoring and itinerary construction be SEPARATE steps, or INTEGRATED?", 16, NAVY, True),
         ("DLIR argues separate → sub-optimal; integration helps. We test this directly.", 14, ORANGE, True)],
        BG_C, ORANGE)
    footer(s, n)


def s_rqs(n):
    s = slide()
    header(s, "Research questions", "scope")
    bullets(s, 0.7, 1.6, 12.0, 5.0, [
        [("RQ1", INDIGO, True),
         (" — Can a personalized, context-aware next-POI model (graph + sequence + user) reach a "
          "competitive tier on a standard check-in benchmark?", GREY, False)],
        [("RQ2", INDIGO, True),
         (" — Can that model be decoded into coherent, loop-free itineraries without dedicated "
          "itinerary training? (the decoupled approach)", GREY, False)],
        [("RQ3", INDIGO, True),
         (" — Does an integrated, end-to-end itinerary model improve over the decoupled decoding — "
          "i.e. does the integration hypothesis hold on our data?", GREY, False)],
        [("RQ4", INDIGO, True),
         (" — Where do our itinerary models stand against the published literature, under the exact "
          "comparable protocol?", GREY, False)],
    ], size=19, gap=18)
    footer(s, n)


def s_contrib(n):
    s = slide()
    header(s, "Contributions", "what's new")
    bullets(s, 0.7, 1.6, 12.0, 5.0, [
        [("A reproducible, honest ", GREY, False), ("personalized context-aware next-POI engine", NAVY, True),
         (" (graph + spatio-temporal context + user embedding) on Foursquare NYC.", GREY, False)],
        [("Two itinerary strategies and a ", GREY, False),
         ("controlled decoupled-vs-integrated comparison", NAVY, True),
         (" (Strategy A: decode the frozen engine; Strategy B: train an integrated pointer).", GREY, False)],
        [("A ", GREY, False), ("literature-comparable Flickr benchmark", NAVY, True),
         (" (Strategy D) with a pairs-F1 harness validated by reproducing Chen 2016.", GREY, False)],
        [("An ", GREY, False), ("honest analysis", NAVY, True),
         (" of where simple baselines suffice and where integration/personalization help — and a "
          "clear path toward a full DLIR-style module.", GREY, False)],
    ], size=18, gap=16)
    footer(s, n)


def s_phase1_arch(n):
    s = slide()
    header(s, "Phase 1 — the personalized, context-aware engine", "next-POI model")
    text(s, 0.55, 1.35, 12.2, 0.5,
         [[("A next-POI model fuses three signals into a score over all POIs:", 16, GREY, False, False)]])
    inputs = [("POI graph\nkNN ∪ co-visit", BG_G, INDIGO), ("context\nΔdistance, Δtime", BG_C, ORANGE),
              ("user\nembedding", BG_U, INDIGO)]
    for i, (t, fill, stroke) in enumerate(inputs):
        x = 0.7 + i * 4.1
        lines = [(ln, 13, stroke if k == 0 else GREY, k == 0) for k, ln in enumerate(t.split("\n"))]
        box(s, x, 2.1, 3.4, 1.15, lines, fill, stroke)
    enc = [("2-layer GCN → POI features", BG_G, INDIGO), ("context MLPs", BG_C, ORANGE),
           ("user vector", BG_U, INDIGO)]
    for i, (t, fill, stroke) in enumerate(enc):
        x = 0.7 + i * 4.1
        arrow(s, x + 1.7, 3.25, x + 1.7, 3.55)
        box(s, x, 3.55, 3.4, 0.8, [(t, 12, stroke, True)], fill, stroke)
    box(s, 0.7, 4.7, 7.0, 0.9, [("concat per step → GRU sequence model → h_T", 14, TEAL, True)], BG_S, TEAL)
    arrow(s, 2.2, 4.35, 2.2, 4.7); arrow(s, 6.3, 4.35, 6.3, 4.7)
    arrow(s, 7.7, 5.15, 8.5, 5.15)
    box(s, 8.5, 4.7, 4.1, 0.9, [("[h_T ; user] → MLP head", 14, ORANGE, True),
                                ("→ score over all POIs", 12, GREY, False)], BG_H, ORANGE)
    arrow(s, 8.8, 4.35, 8.8, 4.7)
    text(s, 0.55, 5.95, 12.2, 0.8,
         [[("This single model carries the thesis's ", 15, GREY, False, False),
           ("“user preferences”", 15, NAVY, True, False), (" (the user embedding) and ", 15, GREY, False, False),
           ("“contextual data”", 15, NAVY, True, False),
           (" (spatial Δd and temporal Δt). It is the engine the itinerary models reuse.", 15, GREY, False, False)]])
    footer(s, n)


def s_phase1_results(n):
    s = slide()
    header(s, "Phase 1 — results (Foursquare NYC, full vocab)", "RQ1")
    rows = [["Metric", "Value"]]
    for k in ["HR@1", "HR@5", "HR@10", "NDCG@10", "MRR"]:
        rows.append([k, f"{PHASE1_METRICS[k]:.3f}"])
    table(s, rows, 0.7, 1.55, 4.2, 2.6, font=14)
    tier = [["Model", "HR@1"]] + [[m, f"{v:.3f}"] for m, v in PHASE1_TIER]
    hl = next(i for i, (m, _v) in enumerate(PHASE1_TIER) if m.startswith("Ours")) + 1
    table(s, tier, 5.4, 1.55, 7.2, 3.6, font=12.5, highlight_row=hl)
    text(s, 0.7, 5.5, 12.0, 1.2,
         [[("HR@1 = 0.187 sits in the ", 16, GREY, False, False), ("LSTM / STGCN tier", 16, NAVY, True, False),
           (" of the LLM4POI benchmark — an honest baseline, below the transformer/LLM SOTA but a strong, "
            "leakage-free engine.", 16, GREY, False, False)],
          [("(RQ1: yes — a competitive personalized engine.)", 13, GREY_MED, False, True)]])
    footer(s, n)


def s_two_strategies(n):
    s = slide()
    header(s, "From engine to itinerary — two strategies", "RQ2 / RQ3")
    box(s, 4.8, 1.55, 3.7, 1.0, [("the next-POI engine", 15, INDIGO, True), ("(Phase 1)", 12, GREY, False)],
        BG_G, INDIGO)
    arrow(s, 5.2, 2.55, 3.5, 3.4); arrow(s, 8.1, 2.55, 9.8, 3.4)
    box(s, 0.7, 3.4, 5.4, 2.2,
        [("Strategy A — DECOUPLED", 16, ORANGE, True), ("decode the FROZEN engine", 13, GREY, False),
         ("roll it out: mask visited,", 12, GREY, False), ("reserve the end POI", 12, GREY, False),
         ("(no itinerary training)", 12, GREY_MED, False)], BG_H, ORANGE)
    box(s, 7.2, 3.4, 5.4, 2.2,
        [("Strategy B — INTEGRATED", 16, TEAL, True), ("TRAIN a pointer end-to-end", 13, GREY, False),
         ("on whole trajectories", 12, GREY, False), ("(GCN + GRU + pointer)", 12, GREY, False),
         ("v1 no context, v2 + context", 12, GREY_MED, False)], BG_S, TEAL)
    text(s, 0.7, 5.85, 12.0, 0.8,
         [[("A = Halder's “separate problems” baseline; B = the “integrated” model. ", 15, GREY, False, False),
           ("Comparing them tests the integration hypothesis directly.", 15, NAVY, True, False)]])
    footer(s, n)


def s_avsb(n):
    s = slide()
    header(s, "The integration experiment — A vs B (NYC, len≥3)", "RQ3 — an honest result")
    rows = [["Method", "pairs-F1", "kind"]]
    order = ["A — frozen rollout (greedy)", "A — frozen rollout (beam 3)",
             "B-v1 — pointer (no context)", "B-v2 — pointer (+ context)"]
    for m in order:
        pf1, _sf, _ex, kind = AVSB_NYC[m]
        rows.append([m, f"{pf1:.3f}", kind])
    table(s, rows, 0.7, 1.55, 8.0, 2.3, font=13, highlight_row=1)
    box(s, 9.0, 1.55, 3.6, 2.3,
        [("Finding", 16, RED, True), ("the INTEGRATED model", 13, GREY, True),
         ("did NOT beat the", 13, GREY, False), ("decoupled rollout", 13, GREY, True),
         ("on NYC.", 13, GREY, False)], BG_H, RED)
    text(s, 0.7, 4.2, 12.0, 2.4,
         [[("Contrary to the integration hypothesis, the trained pointer (B) trailed the frozen rollout (A) "
            "by ~0.03 pairs-F1.", 16, NAVY, True, False)],
          [("Why: the next-POI engine was trained on ~75k prefix→next examples; the integrated pointer saw "
            "only ~10k whole trajectories.", 15, GREY, False, False)],
          [("Lesson: integration does not automatically win — supervision density matters. This nuances "
            "Halder/DLIR rather than contradicting it.", 15, GREY, False, False)]], )
    footer(s, n)


def s_comparability(n):
    s = slide()
    header(s, "But are these numbers comparable?", "the catch")
    text(s, 0.55, 1.35, 12.2, 0.6,
         [[("Our NYC itinerary pairs-F1 (~0.29) looks far below the literature's ~0.6–0.8 — but it's a ", 16, GREY, False, False),
           ("different benchmark", 16, NAVY, True, False), (".", 16, GREY, False, False)]])
    image(s, ASSETS["scale"], 1.5, 2.2, 10.3)
    text(s, 0.55, 6.4, 12.2, 0.7,
         [[("Comparable needs the same dataset + protocol + metric. The literature uses the small ", 15, GREY, False, False),
           ("Flickr", 15, TEAL, True, False), (" datasets → Strategy D.", 15, GREY, False, False)]],
         align=PP_ALIGN.CENTER)
    footer(s, n)


def s_strategyD(n):
    s = slide()
    header(s, "Strategy D — the literature-comparable benchmark", "RQ4")
    bullets(s, 0.7, 1.45, 7.0, 3.5, [
        [("Data: ", NAVY, True), ("5 Flickr photo-trajectory cities (27–88 POIs), Chen 2016's own files → "
                                  "trajectories identical to the literature's.", GREY, False)],
        [("Protocol: ", NAVY, True), ("leave-one-trajectory-out CV, first + last POI given, length ≥ 3.", GREY, False)],
        [("Metric: ", NAVY, True), ("point-F1 and order-aware pairs-F1 (Chen 2016).", GREY, False)],
        [("Methods: ", NAVY, True), ("classical baselines (Random / Popularity / Markov) + the GCN-pointer.", GREY, False)],
    ], size=15.5, gap=12)
    rows = [["City", "#POIs", "#traj≥3"]]
    for c in CITIES:
        p, _u, _t, e = DATASET_STATS[c]
        rows.append([c, p, e])
    table(s, rows, 8.2, 1.6, 4.4, 3.1, font=12.5)
    text(s, 0.7, 5.7, 12.0, 0.8,
         [[("Small vocabularies → pairs-F1 can reach 0.5–0.85 here, vs ~0.29 on NYC's ~5,000 POIs "
            "(the “easier exam”).", 14, GREY_MED, False, True)]])
    footer(s, n)


def s_validation(n):
    s = slide()
    header(s, "Validation — we reproduce Chen 2016", "is the comparison fair?")
    text(s, 0.55, 1.3, 6.6, 1.8,
         [[("A ", 15, GREY, False, False), ("Random", 15, NAVY, True, False),
           (" baseline has no modelling choices — it can only match Chen 2016 if our data, protocol and "
            "metric are identical. They match within noise.", 15, GREY, False, False)]])
    rows = [["pairs-F1", "Tor", "Osa", "Gla", "Edi", "Mel"]]
    for m in ("Random", "PoiPopularity"):
        rows.append([f"{m} ours"] + [f"{OURS_PAIRS_F1[m][c]:.3f}" for c in CITIES])
        rows.append([f"{m} Chen"] + [f"{CHEN_PAIRS_F1[m][c]:.3f}" for c in CITIES])
    table(s, rows, 0.55, 3.5, 6.6, 2.4, font=12)
    image(s, ASSETS["validation"], 7.7, 1.4, 5.2)
    footer(s, n)


def s_results(n):
    s = slide()
    header(s, "Strategy D — results on the literature's scale", "RQ4")
    image(s, ASSETS["results"], 0.6, 1.4, 8.9)
    box(s, 9.75, 1.6, 3.1, 4.3,
        [("Headline", 16, NAVY, True), ("", 6, NAVY, False),
         ("classical", 13, GREY, True), ("0.23 – 0.59", 20, TEAL, True), ("", 5, NAVY, False),
         ("pointer", 13, GREY, True), ("0.31 – 0.49", 20, ORANGE, True), ("", 5, NAVY, False),
         ("vs NYC", 13, GREY, True), ("0.26 – 0.29", 17, RED, True), ("", 5, NAVY, False),
         ("published", 13, GREY, True), ("0.26 – 0.85", 15, NAVY, True)], BG_LIGHT, NAVY)
    text(s, 0.6, 6.45, 9.0, 0.7,
         [[("On the standard benchmark our methods are on the published scale; the simple Markov baseline is "
            "strong and the learned pointer reaches the scale but doesn't beat it.", 14, GREY, False, False)]])
    footer(s, n)


def s_halder(n):
    s = slide()
    header(s, "Positioning vs Halder / DLIR (2024–2025)", "engaging the field")
    box(s, 0.7, 1.5, 5.9, 2.4,
        [("Halder / DLIR", 16, INDIGO, True), ("integrate POI scoring + itinerary", 13, GREY, False),
         ("+ dynamic temporal interest", 12, GREY, False), ("+ queuing time + time budget", 12, GREY, False),
         ("(Transformer + co-visit GCN)", 12, GREY_MED, False)], BG_G, INDIGO)
    box(s, 6.8, 1.5, 5.8, 2.4,
        [("This thesis", 16, TEAL, True), ("tests decoupled (A) vs integrated (B)", 13, GREY, False),
         ("user embedding + Δd/Δt context", 12, GREY, False), ("validated pairs-F1 on Flickr", 12, GREY, False),
         ("(GCN + GRU pointer)", 12, GREY_MED, False)], BG_S, TEAL)
    bullets(s, 0.7, 4.2, 12.0, 2.4, [
        [("We engage Halder's central claim directly", NAVY, True),
         (" — and find integration did not automatically win on our data (a nuance, not a contradiction).", GREY, False)],
        [("We deliberately scope smaller", NAVY, True),
         (": we do not yet model queuing, explicit time budgets, or dynamic time-of-day interest.", GREY, False)],
    ], size=16, gap=12)
    footer(s, n)


def s_limits(n):
    s = slide()
    header(s, "Honest findings & limitations", "self-critique")
    bullets(s, 0.7, 1.6, 12.0, 5.2, [
        [("Integration didn't beat decoupling on NYC", NAVY, True),
         (" — driven by supervision density, not a refutation of the idea.", GREY, False)],
        [("Personalization is uneven", NAVY, True),
         (" — strong on Foursquare (user embedding), weak in the comparable Flickr results.", GREY, False)],
        [("Context is light", NAVY, True),
         (" — spatial Δd + temporal Δt only; no time-of-day dynamics, queuing, or budgets (cf. DLIR).", GREY, False)],
        [("Markov is a hard baseline", NAVY, True),
         (" — on tiny data, a simple transition model beats the neural pointer.", GREY, False)],
        [("Single seed; NYC only for Phase 1", NAVY, True),
         (" — multi-seed and TKY runs are planned.", GREY, False)],
    ], size=17, gap=12)
    footer(s, n)


def s_future(n):
    s = slide()
    header(s, "Future work — toward the full Smart Visit Module", "next")
    bullets(s, 0.7, 1.6, 12.0, 5.2, [
        [("Integrate properly", NAVY, True),
         (" — warm-start / share the engine with the pointer; emit prefix sub-trajectories to match "
          "supervision density.", GREY, False)],
        [("Turn on personalization in the itinerary", NAVY, True),
         (" — the pointer's user-embedding lever; report personalized-vs-not pairs-F1.", GREY, False)],
        [("Richer context", NAVY, True),
         (" — time-of-day / day-of-week dynamic interest, opening hours.", GREY, False)],
        [("Add scheduling", NAVY, True),
         (" — explicit time budget + queuing time (the DLIR direction) for realistic day plans.", GREY, False)],
        [("Robustness", NAVY, True), (" — multi-seed, more cities, ablations of each component.", GREY, False)],
    ], size=17, gap=12)
    footer(s, n)


def s_conclusion(n):
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY, line=None, rounded=False)
    rect(s, 0, 1.7, SW, 0.05, TEAL, line=None, rounded=False)
    text(s, 0.8, 0.6, SW - 1.6, 1.0, [[("Conclusion", 34, WHITE, True, False)]])
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.1), Inches(SW - 1.6), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    pts = [
        "A personalized, context-aware next-POI engine (RQ1) — honest LSTM/STGCN tier on Foursquare NYC.",
        "Two ways to turn it into itineraries (RQ2/RQ3): decoupled decoding vs an integrated pointer — and integration did not automatically win.",
        "A validated, literature-comparable Flickr benchmark (RQ4): Random reproduces Chen 2016; our methods land on the published 0.3–0.85 scale.",
        "Positioned against Halder/DLIR: we test the integration claim and chart the path to a full scheduling-aware module.",
    ]
    for i, t in enumerate(pts):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        b = p.add_run(); b.text = "▸  "; b.font.size = Pt(18); b.font.color.rgb = TEAL; b.font.bold = True
        r = p.add_run(); r.text = t; r.font.size = Pt(18); r.font.color.rgb = WHITE; r.font.name = "Calibri"
    text(s, 0.8, 6.6, SW - 1.6, 0.5, [[("Thank you — questions?", 20, ORANGE, True, False)]])


def s_refs(n):
    s = slide()
    header(s, "Key references", "backup")
    bullets(s, 0.7, 1.6, 12.0, 5.0, [
        "Halder, Lim, Chan, Zhang. Deep Learning of Dynamic POI Generation & Optimisation for Itinerary Rec (DLIR). ACM TORS 2025.",
        "Halder et al. A survey on personalized itinerary recommendation: from optimisation to deep learning. ASOC 2024.",
        "Chen, Ong, Xie. Learning Points and Routes to Recommend Trajectories. CIKM 2016 — pairs-F1 + our Flickr data.",
        "Gao et al. DeepTrip (SIGSPATIAL 2019); SelfTrip (KBS 2022). · Shu et al. AR-Trip (SIGIR 2024).",
        "Yang et al. GETNext (SIGIR 2022); Li et al. LLM4POI (SIGIR 2024) — next-POI tier we compare to.",
        "Krichene & Rendle. On Sampled Metrics for Item Recommendation. KDD 2020 — no cross-protocol comparison.",
    ], size=14, gap=8)
    footer(s, n)


def build():
    s_title()
    s_motivation(2)
    s_gap(3)
    s_rqs(4)
    s_contrib(5)
    s_phase1_arch(6)
    s_phase1_results(7)
    s_two_strategies(8)
    s_avsb(9)
    s_comparability(10)
    s_strategyD(11)
    s_validation(12)
    s_results(13)
    s_halder(14)
    s_limits(15)
    s_future(16)
    s_conclusion(17)
    s_refs(18)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
