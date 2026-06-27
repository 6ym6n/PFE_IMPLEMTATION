"""Model-centric deck (~14 slides): lead with the next-POI engine as THE model,
and Strategy A (decoding it into itineraries) as the itinerary method. Strategy B/D
appear only as supporting experiment/validation. Run:

    py -3.11 presentation/strategy_d/make_model_pptx.py

Output: presentation/strategy_d/model_slides.pptx
"""

from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import charts  # noqa: E402
from flickr_results_data import (  # noqa: E402
    AVSB_NYC, CITIES, OURS_PAIRS_F1, PERSONALIZATION, PHASE1_METRICS, PHASE1_TIER,
    POINTER_PAIRS_F1,
)

ASSETS = charts.make_charts(os.path.join(HERE, "assets"))
OUTPUT = os.path.join(HERE, "model_slides.pptx")

NAVY = RGBColor(0x1B, 0x2A, 0x4E)
INDIGO = RGBColor(0x3F, 0x51, 0xB5)
TEAL = RGBColor(0x0E, 0x7C, 0x7B)
ORANGE = RGBColor(0xE0, 0x7A, 0x1F)
RED = RGBColor(0xC0, 0x39, 0x2B)
GREY = RGBColor(0x55, 0x55, 0x55)
GREY_MED = RGBColor(0x88, 0x88, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_G = RGBColor(0xE3, 0xF2, 0xFD)
BG_C = RGBColor(0xFF, 0xF3, 0xE0)
BG_U = RGBColor(0xF3, 0xE5, 0xF5)
BG_S = RGBColor(0xE8, 0xF5, 0xE9)
BG_H = RGBColor(0xFF, 0xEB, 0xEE)
BG_LIGHT = RGBColor(0xEC, 0xEE, 0xF3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]


def slide():
    return prs.slides.add_slide(BLANK)


def _ns(shp):
    shp.shadow.inherit = False


def rect(s, x, y, w, h, fill, line=None, line_w=1.5, rounded=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    _ns(shp)
    return shp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for pi, para in enumerate(runs):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        for (txt, size, color, bold, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic; r.font.name = "Calibri"
    return tb


def box(s, x, y, w, h, lines, fill, line, anchor=MSO_ANCHOR.MIDDLE):
    shp = rect(s, x, y, w, h, fill, line, line_w=1.75)
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    for i, (txt, size, col, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = col; r.font.name = "Calibri"
    return shp


def arrow(s, x1, y1, x2, y2, color=NAVY, width=2.0):
    conn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color; conn.line.width = Pt(width)
    ln = conn.line._get_or_add_ln()
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    _ns(conn)
    return conn


def header(s, title, kicker=None):
    rect(s, 0, 0, SW, 1.15, NAVY, line=None, rounded=False)
    rect(s, 0, 1.15, SW, 0.06, TEAL, line=None, rounded=False)
    text(s, 0.55, 0.18, SW - 1.1, 0.85, [[(title, 27, WHITE, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(s, 0.55, 0.04, SW - 1.1, 0.3, [[(kicker, 11, RGBColor(0xBF, 0xC7, 0xD9), False, True)]])


def footer(s, n):
    text(s, SW - 1.4, SH - 0.45, 1.1, 0.3, [[(f"{n}", 11, GREY_MED, False, False)]], align=PP_ALIGN.RIGHT)
    text(s, 0.55, SH - 0.45, 10, 0.3,
         [[("A personalized, context-aware itinerary recommender", 9, GREY_MED, False, True)]])


def bullets(s, x, y, w, h, items, size=17, gap=10):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.space_before = Pt(2)
        if isinstance(item, str):
            item = [(item, GREY, False)]
        b = p.add_run(); b.text = "▸  "; b.font.size = Pt(size); b.font.color.rgb = TEAL; b.font.bold = True
        for (txt, color, bold) in item:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = "Calibri"
    return tb


def image(s, path, x, y, w):
    return s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


def table(s, rows, x, y, w, h, highlight_row=None, font=12):
    nr, nc = len(rows), len(rows[0])
    gt = s.shapes.add_table(nr, nc, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for ci in range(nc):
        for ri in range(nr):
            cell = gt.cell(ri, ci)
            cell.margin_top = Pt(3); cell.margin_bottom = Pt(3); cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(rows[ri][ci]); r.font.size = Pt(font); r.font.name = "Calibri"
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY; r.font.color.rgb = WHITE; r.font.bold = True
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE if ri % 2 else BG_LIGHT
                r.font.color.rgb = GREY
                if highlight_row is not None and ri == highlight_row:
                    cell.fill.fore_color.rgb = TEAL; r.font.color.rgb = WHITE; r.font.bold = True
    return gt


# ===========================================================================
# Slides
# ===========================================================================
def s_title():
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY, line=None, rounded=False)
    rect(s, 0, 3.3, SW, 0.05, TEAL, line=None, rounded=False)
    text(s, 1.0, 1.6, SW - 2, 1.6,
         [[("The Model", 40, WHITE, True, False)],
          [("A Personalized, Context-Aware Itinerary Recommender", 24, WHITE, True, False)]],
         align=PP_ALIGN.CENTER)
    text(s, 1.0, 3.5, SW - 2, 0.7,
         [[("A next-POI engine — GCN + GRU + user embedding + context", 18, RGBColor(0x9F, 0xC0, 0xC0), False, True)]],
         align=PP_ALIGN.CENTER)
    text(s, 1.0, 4.6, SW - 2, 0.9,
         [[("…decoded into full tourist itineraries (Frozen Rollout)", 18, WHITE, False, False)]],
         align=PP_ALIGN.CENTER)
    text(s, 1.0, 6.4, SW - 2, 0.5,
         [[("Master's thesis · [Your name] · [date]", 13, RGBColor(0xBF, 0xC7, 0xD9), False, False)]],
         align=PP_ALIGN.CENTER)


def s_goal(n):
    s = slide()
    header(s, "The goal", "problem")
    bullets(s, 0.7, 1.6, 12.0, 2.6, [
        [("A tourist with limited time wants a ", GREY, False), ("personalized day plan", NAVY, True),
         (" — an ", GREY, False), ("ordered route", NAVY, True), (" of places to visit.", GREY, False)],
        [("The plan must reflect ", GREY, False), ("who they are", NAVY, True),
         (" (preferences) and ", GREY, False), ("context", NAVY, True),
         (" (distances, time).", GREY, False)],
    ], size=18)
    box(s, 1.6, 4.6, 10.1, 1.4,
        [("My approach: one model does both —", 18, NAVY, True),
         ("a personalized next-POI engine, decoded into an itinerary.", 18, NAVY, True)], BG_S, TEAL)
    footer(s, n)


def s_overview(n):
    s = slide()
    header(s, "The model in one picture", "overview")
    box(s, 0.7, 2.7, 4.0, 1.7,
        [("THE MODEL", 13, INDIGO, True), ("next-POI engine", 16, INDIGO, True),
         ("GCN + GRU + user + context", 11, GREY, False)], BG_G, INDIGO)
    arrow(s, 4.7, 3.55, 5.5, 3.55)
    box(s, 5.5, 2.7, 3.6, 1.7,
        [("Frozen Rollout", 13, ORANGE, True), ("roll it out", 16, ORANGE, True),
         ("(decode the engine)", 11, GREY, False)], BG_H, ORANGE)
    arrow(s, 9.1, 3.55, 9.9, 3.55)
    box(s, 9.9, 2.7, 2.7, 1.7,
        [("itinerary", 16, TEAL, True), ("ordered route", 12, GREY, False)], BG_S, TEAL)
    text(s, 0.7, 5.0, 12.0, 1.2,
         [[("My ", 18, GREY, False, False), ("model", 18, INDIGO, True, False),
           (" is a next-POI engine; my ", 18, GREY, False, False),
           ("itinerary method", 18, ORANGE, True, False),
           (" decodes it. Everything else (training a separate model, the benchmark) is an "
            "experiment around this.", 18, GREY, False, False)]])
    footer(s, n)


def s_arch(n):
    s = slide()
    header(s, "The model — architecture", "THE model")
    inputs = [("POI graph\nkNN ∪ co-visit", BG_G, INDIGO), ("context\nΔdistance, Δtime", BG_C, ORANGE),
              ("user\nidentity", BG_U, INDIGO)]
    for i, (t, fill, stroke) in enumerate(inputs):
        x = 0.7 + i * 4.1
        lines = [(ln, 13, stroke if k == 0 else GREY, k == 0) for k, ln in enumerate(t.split("\n"))]
        box(s, x, 1.5, 3.4, 1.0, lines, fill, stroke)
    enc = [("2-layer GCN", BG_G, INDIGO), ("context MLPs", BG_C, ORANGE), ("user embedding", BG_U, INDIGO)]
    for i, (t, fill, stroke) in enumerate(enc):
        x = 0.7 + i * 4.1
        arrow(s, x + 1.7, 2.5, x + 1.7, 2.8)
        box(s, x, 2.8, 3.4, 0.75, [(t, 13, stroke, True)], fill, stroke)
    box(s, 0.7, 3.95, 7.0, 0.85, [("concat per step → GRU → h_T", 14, TEAL, True)], BG_S, TEAL)
    arrow(s, 2.2, 3.55, 2.2, 3.95); arrow(s, 6.3, 3.55, 6.3, 3.95)
    arrow(s, 7.7, 4.37, 8.6, 4.37)
    box(s, 8.6, 3.95, 4.0, 0.85, [("[h_T ; user] → MLP head", 13, ORANGE, True)], BG_H, ORANGE)
    arrow(s, 9.0, 3.55, 9.0, 3.95)
    box(s, 3.0, 5.25, 7.3, 0.85, [("→ a score for every POI (what to visit next)", 15, NAVY, True)], BG_LIGHT, NAVY)
    arrow(s, 6.6, 4.8, 6.6, 5.25)
    text(s, 0.7, 6.4, 12.0, 0.6,
         [[("Three signals → graph + sequence + user → a next-POI score. ", 14, GREY, False, False),
           ("This is the model I built.", 14, NAVY, True, False)]])
    footer(s, n)


def s_parts(n):
    s = slide()
    header(s, "How each part works", "the model, explained")
    bullets(s, 0.7, 1.55, 12.0, 5.2, [
        [("POI graph + GCN", INDIGO, True),
         (" — places are linked if they're close or often visited together; the GCN gives each POI a "
          "“fingerprint” aware of its neighbourhood.", GREY, False)],
        [("Context (Δd, Δt)", ORANGE, True),
         (" — how far and how long since the last place; small networks turn these into vectors.", GREY, False)],
        [("User embedding", INDIGO, True),
         (" — a small “taste profile” per tourist → the title's user preferences.", GREY, False)],
        [("GRU sequence model", TEAL, True),
         (" — reads the places visited so far and summarises them into one state.", GREY, False)],
        [("Scoring head", ORANGE, True),
         (" — combines the state + the user, and scores every POI for “what to visit next”.", GREY, False)],
    ], size=16.5, gap=13)
    footer(s, n)


def s_title_fit(n):
    s = slide()
    header(s, "The model carries the thesis title", "why this model")
    text(s, 0.7, 1.5, 12.0, 0.6,
         [[("“Personalized Tourist Itinerary Recommendation Based on ", 17, GREY, False, True),
           ("User Preferences", 17, INDIGO, True, True), (" and ", 17, GREY, False, True),
           ("Contextual Data", 17, ORANGE, True, True), ("”", 17, GREY, False, True)]])
    box(s, 0.9, 2.7, 5.6, 2.4,
        [("User preferences", 18, INDIGO, True), ("", 6, INDIGO, False),
         ("= the user embedding", 15, GREY, True), ("a learned taste profile", 13, GREY, False),
         ("for each tourist", 13, GREY, False)], BG_U, INDIGO)
    box(s, 6.9, 2.7, 5.6, 2.4,
        [("Contextual data", 18, ORANGE, True), ("", 6, ORANGE, False),
         ("= Δdistance and Δtime", 15, GREY, True), ("spatial + temporal", 13, GREY, False),
         ("gaps between places", 13, GREY, False)], BG_C, ORANGE)
    text(s, 0.7, 5.6, 12.0, 0.8,
         [[("Both are ", 16, GREY, False, False), ("inside the model", 16, NAVY, True, False),
           (" — so the title is a description of the architecture, not an aspiration.", 16, GREY, False, False)]])
    footer(s, n)


def s_engine_results(n):
    s = slide()
    header(s, "The engine works (Phase 1)", "evidence the model is sound")
    rows = [["Metric", "Value"]]
    for k in ["HR@1", "HR@5", "HR@10", "MRR"]:
        rows.append([k, f"{PHASE1_METRICS[k]:.3f}"])
    table(s, rows, 0.7, 1.55, 4.0, 2.2, font=14)
    tier = [["Model", "HR@1"]] + [[m, f"{v:.3f}"] for m, v in PHASE1_TIER]
    hl = next(i for i, (m, _v) in enumerate(PHASE1_TIER) if m.startswith("Ours")) + 1
    table(s, tier, 5.2, 1.55, 7.4, 3.6, font=12.5, highlight_row=hl)
    text(s, 0.7, 5.5, 12.0, 1.0,
         [[("On Foursquare NYC the engine reaches ", 16, GREY, False, False),
           ("HR@1 = 0.187", 16, NAVY, True, False),
           (" — the honest LSTM/STGCN tier. A strong, leakage-free engine to build itineraries on.",
            16, GREY, False, False)]])
    footer(s, n)


def s_strategyA(n):
    s = slide()
    header(s, "Frozen Rollout — turning the engine into an itinerary", "my itinerary method")
    text(s, 0.55, 1.35, 12.2, 0.5,
         [[("Roll the engine out: build the route one POI at a time.", 16, GREY, False, False)]])
    steps = [("route = [start]", BG_U, INDIGO), ("engine scores\nall POIs", BG_G, INDIGO),
             ("mask visited;\nreserve the end", BG_C, ORANGE), ("append the\nbest POI", BG_H, ORANGE)]
    for i, (t, fill, stroke) in enumerate(steps):
        x = 0.6 + i * 3.15
        lines = [(ln, 13, stroke if k == 0 else GREY, k == 0) for k, ln in enumerate(t.split("\n"))]
        box(s, x, 2.3, 2.7, 1.3, lines, fill, stroke)
        if i < 3:
            arrow(s, x + 2.7, 2.95, x + 2.7 + 0.45, 2.95)
    # loop-back
    arrow(s, 11.05, 3.6, 11.05, 4.1, color=TEAL)
    arrow(s, 11.05, 4.1, 1.95, 4.1, color=TEAL)
    arrow(s, 1.95, 4.1, 1.95, 3.6, color=TEAL)
    text(s, 4.0, 4.15, 5.0, 0.4, [[("repeat until length K, ending at e", 12, TEAL, True, True)]],
         align=PP_ALIGN.CENTER)
    bullets(s, 0.7, 5.0, 12.0, 1.8, [
        [("No new training", NAVY, True), (" — reuse the frozen engine; the itinerary is an inference-time decoder.", GREY, False)],
        [("Always valid", NAVY, True), (" — loop-free, exactly K stops, starts at s and ends at e.", GREY, False)],
    ], size=15)
    footer(s, n)


def s_example(n):
    s = slide()
    header(s, "Frozen Rollout — worked example (Edinburgh)", "how it builds a route")
    box(s, 0.7, 1.55, 4.2, 1.2, [("QUERY", 14, INDIGO, True), ("start 15 · end 16 · K = 4", 14, GREY, False)],
        BG_U, INDIGO)
    seqs = [
        ("step 1", ["15"], "engine → best next = 3"),
        ("step 2", ["15", "3"], "engine → best next = 13"),
        ("step 3", ["15", "3", "13"], "last hop → reserved end 16"),
        ("result", ["15", "3", "13", "16"], "matches the real trip ✓"),
    ]
    for i, (label, seq, note) in enumerate(seqs):
        y = 3.1 + i * 0.95
        text(s, 0.7, y, 1.3, 0.6, [[(label, 12, NAVY, True, False)]], anchor=MSO_ANCHOR.MIDDLE)
        for j, p in enumerate(seq):
            last = (i == 3)
            box(s, 2.1 + j * 0.95, y, 0.8, 0.6, [(p, 14, TEAL if last else NAVY, True)],
                BG_S if last else BG_LIGHT, TEAL if last else NAVY)
        text(s, 6.4, y, 6.0, 0.6, [[(note, 13, GREY, False, True)]], anchor=MSO_ANCHOR.MIDDLE)
    footer(s, n)


def s_A_results(n):
    s = slide()
    header(s, "Frozen Rollout — results, and why it's the headline", "my result")
    rows = [["Method (NYC, len≥3)", "pairs-F1", "kind"]]
    _relabel = {"A — frozen rollout (greedy)": "Frozen Rollout (greedy)",
                "A — frozen rollout (beam 3)": "Frozen Rollout (beam 3)",
                "B-v1 — pointer (no context)": "Trained Pointer (no context)",
                "B-v2 — pointer (+ context)": "Trained Pointer (+ context)"}
    for m in _relabel:
        pf1, _s, _e, kind = AVSB_NYC[m]
        rows.append([_relabel[m], f"{pf1:.3f}", kind])
    table(s, rows, 0.7, 1.55, 8.0, 2.3, font=13, highlight_row=1)
    box(s, 9.0, 1.55, 3.6, 2.3,
        [("Frozen Rollout wins", 16, TEAL, True), ("decoding the strong", 13, GREY, False),
         ("engine beats training", 13, GREY, False), ("a separate model.", 13, GREY, False)], BG_S, TEAL)
    bullets(s, 0.7, 4.2, 12.0, 2.2, [
        [("Best of my itinerary methods", NAVY, True), (" on NYC (~0.29 pairs-F1).", GREY, False)],
        [("Simplest and most interpretable", NAVY, True), (" — easy to explain and trust.", GREY, False)],
        [("Reuses the personalized engine directly", NAVY, True),
         (" — so the itinerary inherits user preferences + context.", GREY, False)],
    ], size=16, gap=11)
    footer(s, n)


def s_validation(n):
    s = slide()
    header(s, "Validation — my approach is on the literature's scale", "rigor")
    image(s, ASSETS["results"], 0.6, 1.45, 8.7)
    bullets(s, 9.5, 1.7, 3.5, 4.5, [
        [("Reproduces Chen 2016", NAVY, True), (" (Random baseline) → the setup is fair.", GREY, False)],
        [("On the published", NAVY, True), (" 0.3–0.85 pairs-F1 scale.", GREY, False)],
        [("Honest:", RED, True), (" a simple Markov baseline is strong on tiny data.", GREY, False)],
    ], size=14, gap=12)
    text(s, 0.6, 6.5, 12.0, 0.6,
         [[("The Flickr Benchmark re-runs everything on the standard Flickr datasets — validating, not replacing, the "
            "headline model.", 14, GREY, False, True)]])
    footer(s, n)


def s_beyond(n):
    s = slide()
    header(s, "What else I tested (and found)", "scientific rigor")
    bullets(s, 0.7, 1.6, 12.0, 4.4, [
        [("Does an integrated model beat decoupling?", NAVY, True),
         (" I trained a dedicated pointer (the Trained Pointer) end-to-end. It did NOT beat the Frozen Rollout — because the "
          "engine has far denser training signal. (Tests Halder/DLIR's integration claim.)", GREY, False)],
        [("Is the neural model better than simple baselines?", NAVY, True),
         (" On the small Flickr data, a Markov transition baseline is the strongest of my methods — an honest "
          "result I quantified rather than hid.", GREY, False)],
        [("Does personalization help?", NAVY, True),
         (" I ablate the user embedding (on/off) per city — measured on the next slide.", GREY, False)],
    ], size=16.5, gap=15)
    footer(s, n)


def s_personalization(n):
    s = slide()
    header(s, "Does personalization help? (measured)", "the “user preferences” claim")
    text(s, 0.55, 1.35, 12.2, 0.5,
         [[("I switch the user embedding OFF vs ON (identical otherwise) on the comparable Flickr benchmark:",
            15, GREY, False, False)]])
    rec = {"Osaka": "28%", "Glasgow": "32%", "Toronto": "55%"}
    rows = [["City", "no user", "with user", "delta", "recurring users"]]
    for c in ("Osaka", "Glasgow", "Toronto"):
        d = PERSONALIZATION[c]
        rows.append([c, f"{d['no_user']:.3f}", f"{d['user']:.3f}", f"{d['user'] - d['no_user']:+.3f}", rec[c]])
    rows.append(["Edinburgh, Melbourne", "—", "—", "(GPU run)", "59%, 61%"])
    table(s, rows, 1.6, 2.1, 10.1, 2.6, font=13, highlight_row=2)
    bullets(s, 0.7, 5.1, 12.0, 1.8, [
        [("Mixed / near-neutral", NAVY, True),
         (" — a clear gain on Glasgow (+0.038), roughly neutral on Osaka and Toronto.", GREY, False)],
        [("Honest reason", NAVY, True),
         (" — in leave-one-out, even recurring users have very few trips, so the embedding is weakly "
          "estimated (cold-start). Personalization is a real lever, not a guaranteed win on small data.",
          GREY, False)],
    ], size=15)
    footer(s, n)


def s_future(n):
    s = slide()
    header(s, "Limitations & future work", "honesty")
    bullets(s, 0.7, 1.6, 12.0, 5.0, [
        [("Light context", NAVY, True), (" — Δd/Δt only; add time-of-day, opening hours, queuing.", GREY, False)],
        [("No scheduling", NAVY, True), (" — add explicit time budgets (the DLIR direction) for realistic days.", GREY, False)],
        [("Integrate properly", NAVY, True), (" — share/warm-start the engine with the trained model to lift it past A.", GREY, False)],
        [("Stronger personalization", NAVY, True), (" — richer user features beyond a single embedding.", GREY, False)],
        [("Robustness", NAVY, True), (" — multi-seed, more cities, component ablations.", GREY, False)],
    ], size=17, gap=12)
    footer(s, n)


def s_conclusion(n):
    s = slide()
    rect(s, 0, 0, SW, SH, NAVY, line=None, rounded=False)
    rect(s, 0, 1.7, SW, 0.05, TEAL, line=None, rounded=False)
    text(s, 0.8, 0.6, SW - 1.6, 1.0, [[("My model — in one sentence", 32, WHITE, True, False)]])
    box(s, 1.2, 2.6, SW - 2.4, 1.9,
        [("A personalized, context-aware next-POI engine", 22, WHITE, True),
         ("(GCN + GRU + user + context), decoded into tourist", 22, WHITE, True),
         ("itineraries by rolling it out (the Frozen Rollout).", 22, WHITE, True)], NAVY, TEAL)
    tb = s.shapes.add_textbox(Inches(0.8), Inches(4.9), Inches(SW - 1.6), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    for i, t in enumerate([
        "Strongest of my itinerary methods, and the simplest to explain.",
        "Built on a validated engine; benchmarked honestly against the literature.",
        "Integration (B) and the Flickr benchmark (D) are experiments that support — not replace — it.",
    ]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        b = p.add_run(); b.text = "▸  "; b.font.size = Pt(16); b.font.color.rgb = TEAL; b.font.bold = True
        r = p.add_run(); r.text = t; r.font.size = Pt(16); r.font.color.rgb = WHITE; r.font.name = "Calibri"
    text(s, 0.8, 6.7, SW - 1.6, 0.5, [[("Thank you — questions?", 18, ORANGE, True, False)]])


def build():
    s_title()
    s_goal(2)
    s_overview(3)
    s_arch(4)
    s_parts(5)
    s_title_fit(6)
    s_engine_results(7)
    s_strategyA(8)
    s_example(9)
    s_A_results(10)
    s_validation(11)
    s_beyond(12)
    s_personalization(13)
    s_future(14)
    s_conclusion(15)
    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    build()
